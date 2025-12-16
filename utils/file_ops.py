"""
文件操作工具模块
提供文件扫描、文本提取等通用功能
支持智能 OCR fallback 机制
"""

import re
import shutil
import time
import logging
import sqlite3
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional

# 配置日志
logger = logging.getLogger(__name__)


# ==================== OCR 缓存系统 ====================


def compute_file_hash(path: Path) -> str:
    """
    计算文件 hash（快速算法：文件大小 + 头部 8KB）
    
    与 Memory 系统保持一致
    
    Args:
        path: 文件路径
    
    Returns:
        SHA256 哈希字符串
    """
    hasher = hashlib.sha256()
    file_size = path.stat().st_size
    
    # 1. 写入文件大小
    hasher.update(str(file_size).encode('utf-8'))
    
    # 2. 读取头部 8KB
    try:
        with open(path, 'rb') as f:
            chunk = f.read(8192)
            hasher.update(chunk)
    except Exception:
        pass
    
    return hasher.hexdigest()


def calculate_quality_score(text: str, confidence: float = 0.0) -> Tuple[int, bool]:
    """
    计算 OCR 结果质量评分
    
    Args:
        text: OCR 提取的文本
        confidence: OCR 置信度（0-1）
    
    Returns:
        (quality_score, needs_review)
        - quality_score: 0-100
        - needs_review: 是否需要人工审查
    
    评分规则：
    - 初始 100 分
    - 文本长度 < 50 扣 30 分，< 100 扣 15 分
    - 乱码率 > 50% 扣 50 分，> 30% 扣 30 分，> 10% 扣 10 分
    - 置信度 < 0.5 扣 20 分，< 0.7 扣 10 分
    - 分数 < 60 标记为需要审查
    """
    score = 100
    text_stripped = text.strip()
    char_count = len(text_stripped)
    
    # 1. 文本长度检查
    if char_count == 0:
        score = 0
    elif char_count < 50:
        score -= 30
    elif char_count < 100:
        score -= 15
    
    # 2. 乱码检查（复用 should_use_ocr 的逻辑）
    if char_count > 0:
        normal_chars = sum(1 for c in text_stripped 
                          if c.isalnum() or c.isspace() or c in '，。！？；：、""''（）【】《》,.!?;:\'"()[]{}')
        garbage_ratio = 1 - (normal_chars / char_count)
        
        if garbage_ratio > 0.5:
            score -= 50
        elif garbage_ratio > 0.3:
            score -= 30
        elif garbage_ratio > 0.1:
            score -= 10
    
    # 3. 置信度影响
    if confidence > 0:
        if confidence < 0.5:
            score -= 20
        elif confidence < 0.7:
            score -= 10
    
    # 4. 确保分数在 0-100 范围内
    score = max(0, min(100, score))
    
    # 5. 判断是否需要审查
    needs_review = score < 60
    
    return score, needs_review


class OCRCache:
    """
    OCR 结果缓存（SQLite）
    
    用于避免重复 OCR 处理，提升性能和降低成本
    """
    
    def __init__(self, db_path: Optional[Path] = None):
        """
        初始化缓存数据库
        
        Args:
            db_path: 数据库文件路径，默认为 ~/.digital_janitor/ocr_cache.db
        """
        if db_path is None:
            # 默认位置：~/.digital_janitor/ocr_cache.db
            cache_dir = Path.home() / ".digital_janitor"
            cache_dir.mkdir(parents=True, exist_ok=True)
            db_path = cache_dir / "ocr_cache.db"
        
        self.db_path = db_path
        self._init_db()
    
    def _init_db(self):
        """初始化数据库表"""
        conn = sqlite3.connect(str(self.db_path))
        cursor = conn.cursor()
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS ocr_cache (
                file_hash TEXT PRIMARY KEY,
                text TEXT NOT NULL,
                method TEXT NOT NULL,
                confidence REAL,
                quality_score INTEGER,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        conn.commit()
        conn.close()
    
    def get(self, file_hash: str) -> Optional[Dict[str, Any]]:
        """
        查询缓存
        
        Args:
            file_hash: 文件 hash
        
        Returns:
            缓存数据字典或 None
        """
        conn = sqlite3.connect(str(self.db_path))
        cursor = conn.cursor()
        cursor.execute(
            "SELECT text, method, confidence, quality_score FROM ocr_cache WHERE file_hash = ?",
            (file_hash,)
        )
        row = cursor.fetchone()
        conn.close()
        
        if row:
            return {
                "text": row[0],
                "method": row[1],
                "confidence": row[2],
                "quality_score": row[3]
            }
        return None
    
    def set(self, file_hash: str, text: str, method: str, confidence: float, quality_score: int):
        """
        保存缓存
        
        Args:
            file_hash: 文件 hash
            text: 提取的文本
            method: 提取方法
            confidence: 置信度
            quality_score: 质量评分
        """
        conn = sqlite3.connect(str(self.db_path))
        cursor = conn.cursor()
        cursor.execute(
            """INSERT OR REPLACE INTO ocr_cache 
               (file_hash, text, method, confidence, quality_score) 
               VALUES (?, ?, ?, ?, ?)""",
            (file_hash, text, method, confidence, quality_score)
        )
        conn.commit()
        conn.close()
    
    def clear(self):
        """清空缓存"""
        conn = sqlite3.connect(str(self.db_path))
        cursor = conn.cursor()
        cursor.execute("DELETE FROM ocr_cache")
        conn.commit()
        conn.close()


def discover_files(inbox: Path, ignore_hidden: bool = True) -> List[Path]:
    """
    扫描目录中的文件
    
    Args:
        inbox: 要扫描的目录
        ignore_hidden: 是否忽略隐藏文件（以 . 开头）
    
    Returns:
        文件路径列表（按文件名排序）
    """
    files: List[Path] = []
    for p in inbox.iterdir():
        if p.is_file():
            if ignore_hidden and p.name.startswith("."):
                continue
            files.append(p)
    return sorted(files, key=lambda x: x.name.lower())


# ==================== 智能 OCR 增强功能 ====================


def should_use_ocr(extracted_text: str, page_count: int) -> Tuple[bool, str]:
    """
    判断是否需要 OCR，返回 (是否需要, 原因)
    
    触发条件：
    1. 完全空文本
    2. 平均每页 < 100 字符
    3. 乱码比例 > 30%
    4. 几乎全是空白符
    
    Args:
        extracted_text: pypdf 提取的文本
        page_count: PDF 页数
    
    Returns:
        (是否需要 OCR, 原因)
    """
    # 条件 1: 完全空文本
    if not extracted_text or len(extracted_text.strip()) == 0:
        return True, "文本为空，可能是扫描件"
    
    text_stripped = extracted_text.strip()
    char_count = len(text_stripped)
    
    # 条件 2: 字符密度过低（平均每页 < 100 字符）
    if page_count > 0:
        chars_per_page = char_count / page_count
        if chars_per_page < 100:
            return True, f"字符密度过低（{chars_per_page:.1f} 字符/页），可能是扫描件"
    
    # 条件 3: 几乎全是空白符（空白符占比 > 90%）
    whitespace_count = sum(1 for c in extracted_text if c.isspace())
    if char_count > 0:
        whitespace_ratio = whitespace_count / len(extracted_text)
        if whitespace_ratio > 0.9:
            return True, f"空白符占比过高（{whitespace_ratio:.1%}），可能是扫描件"
    
    # 条件 4: 乱码比例过高（非常见字符 > 30%）
    # 常见字符：中文、英文、数字、常用标点
    normal_chars = sum(1 for c in text_stripped 
                      if c.isalnum() or c.isspace() or c in '，。！？；：、""''（）【】《》,.!?;:\'"()[]{}')
    if char_count > 0:
        garbage_ratio = 1 - (normal_chars / char_count)
        if garbage_ratio > 0.3:
            return True, f"乱码比例过高（{garbage_ratio:.1%}），可能是扫描件或编码问题"
    
    # 不需要 OCR
    return False, "文本提取正常"


def is_important_document(pdf_path: Path) -> bool:
    """
    判断是否为重要文档（值得用 Vision LLM）
    
    判断依据：
    - 文件名包含关键词：invoice/发票/contract/合同/报告/report
    - 文件大小在 100KB - 10MB 之间
    - 页数 <= 5 页（需要打开 PDF 检查）
    
    Returns:
        True 则优先使用 Vision LLM
    """
    try:
        from config.ocr_config import OCR_CONFIG
        
        # 1. 检查文件名关键词
        filename_lower = pdf_path.name.lower()
        has_keyword = any(kw.lower() in filename_lower for kw in OCR_CONFIG.important_keywords)
        
        # 2. 检查文件大小
        file_size_kb = pdf_path.stat().st_size / 1024
        size_ok = OCR_CONFIG.important_min_size_kb <= file_size_kb <= OCR_CONFIG.important_max_size_kb
        
        # 3. 检查页数（需要打开 PDF）
        page_count = 1  # 默认假设 1 页
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(pdf_path))
            page_count = len(reader.pages)
        except Exception:
            pass
        
        pages_ok = page_count <= OCR_CONFIG.important_max_pages
        
        # 综合判断
        is_important = has_keyword and size_ok and pages_ok
        
        if is_important:
            logger.info(
                f"文档被标记为重要: {pdf_path.name} "
                f"(关键词匹配={has_keyword}, 大小={file_size_kb:.1f}KB, 页数={page_count})"
            )
        
        return is_important
        
    except Exception as e:
        logger.warning(f"判断文档重要性时出错: {e}")
        return False


def extract_with_rapidocr(
    file_path: Path, 
    max_pages: int = 10,
    min_confidence: float = 0.5
) -> Dict[str, Any]:
    """
    使用 RapidOCR 提取文本（支持 PDF 和图片）
    
    Args:
        file_path: PDF 或图片文件路径
        max_pages: 最多处理前 N 页（仅 PDF）
        min_confidence: 最小置信度阈值
    
    Returns:
        {
            "text": str,
            "confidence": float,
            "pages_processed": int,
            "elapsed_ms": int,
            "error": str | None
        }
    """
    start_time = time.time()
    result: Dict[str, Any] = {
        "text": "",
        "confidence": 0.0,
        "pages_processed": 0,
        "elapsed_ms": 0,
        "error": None
    }
    
    try:
        # 1. 导入依赖
        from rapidocr_onnxruntime import RapidOCR
        from config.ocr_config import OCR_CONFIG
        
        # 2. 初始化 RapidOCR
        ocr_engine = RapidOCR()
        
        # 3. 根据文件类型加载图片
        ext = file_path.suffix.lower()
        
        # 🆕 图片文件：直接加载
        if ext in ['.png', '.jpg', '.jpeg', '.webp']:
            try:
                from PIL import Image
                images = [Image.open(str(file_path))]
                logger.info(f"加载图片文件: {file_path.name}")
            except Exception as e:
                result["error"] = f"图片加载失败: {str(e)}"
                return result
        
        # PDF 文件：转换为图片
        elif ext == '.pdf':
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(
                    str(file_path),
                    dpi=OCR_CONFIG.rapidocr_dpi,
                    first_page=1,
                    last_page=min(max_pages, 999)  # 限制最大页数
                )
                logger.info(f"PDF 转图片: {file_path.name}, 页数={len(images)}")
            except Exception as e:
                result["error"] = f"PDF 转图片失败: {str(e)}"
                return result
        
        else:
            result["error"] = f"不支持的文件类型: {ext}"
            return result
        
        # 4. 逐页/逐图 OCR
        all_text = []
        all_confidences = []
        
        for page_num, image in enumerate(images, 1):
            try:
                # RapidOCR 识别
                ocr_result, elapse = ocr_engine(image)
                
                if ocr_result:
                    for line in ocr_result:
                        # line 格式: [box, text, confidence]
                        text = line[1]
                        confidence = line[2]
                        
                        # 过滤低置信度结果
                        if confidence >= min_confidence:
                            all_text.append(text)
                            all_confidences.append(confidence)
                
                result["pages_processed"] = page_num
                
            except Exception as e:
                logger.warning(f"OCR 第 {page_num} 页失败: {e}")
                continue
        
        # 5. 汇总结果
        result["text"] = "\n".join(all_text)
        result["confidence"] = sum(all_confidences) / len(all_confidences) if all_confidences else 0.0
        
    except ImportError as e:
        result["error"] = f"OCR 库未安装: {str(e)}"
        logger.error(result["error"])
    except Exception as e:
        result["error"] = f"OCR 处理失败: {str(e)}"
        logger.error(result["error"])
    finally:
        # 记录耗时
        result["elapsed_ms"] = int((time.time() - start_time) * 1000)
    
    return result


def extract_text_preview_enhanced(path: Path, limit: int = 1000) -> Dict[str, Any]:
    """
    智能文本提取（支持 OCR fallback + 缓存 + 图片）
    
    Args:
        path: 文件路径
        limit: 最大字符数
    
    Returns:
        {
            "text": str,
            "method": "direct" | "rapidocr" | "vision_llm" | "direct_fallback" | "cached",
            "confidence": float,
            "quality_score": int,  # 🆕 质量评分 (0-100)
            "needs_review": bool,  # 🆕 是否需要人工审查
            "page_count": int,
            "char_count": int,
            "processing_time_ms": int,
            "error": str | None
        }
    
    处理流程：
    1. 计算文件 hash，查询缓存
    2. 如果是图片文件，直接使用 RapidOCR
    3. 如果是 PDF，使用 pypdf 尝试标准提取
    4. 调用 should_use_ocr() 判断是否需要 OCR
    5. 计算质量评分，保存缓存
    """
    start_time = time.time()
    ext = path.suffix.lower()
    
    result: Dict[str, Any] = {
        "text": "",
        "method": "unknown",
        "confidence": 0.0,
        "quality_score": 0,
        "needs_review": True,
        "page_count": 0,
        "char_count": 0,
        "processing_time_ms": 0,
        "error": None
    }
    
    # 🆕 初始化缓存
    cache = OCRCache()
    
    # 🆕 计算文件 hash
    try:
        file_hash = compute_file_hash(path)
    except Exception as e:
        logger.warning(f"计算文件 hash 失败: {e}")
        file_hash = None
    
    # 🆕 尝试从缓存加载（仅对 OCR 需求高的文件类型）
    if file_hash and ext in ['.pdf', '.png', '.jpg', '.jpeg', '.webp']:
        cached = cache.get(file_hash)
        if cached:
            logger.info(f"✅ 命中 OCR 缓存: {path.name}")
            result["text"] = cached["text"][:limit]
            result["method"] = f"{cached['method']}_cached"
            result["confidence"] = cached["confidence"]
            result["quality_score"] = cached["quality_score"]
            result["needs_review"] = cached["quality_score"] < 60
            result["char_count"] = len(result["text"])
            result["processing_time_ms"] = 0  # 标记为缓存
            return result
    
    # 🆕 图片文件处理 - 直接使用 OCR
    if ext in ['.png', '.jpg', '.jpeg', '.webp']:
        logger.info(f"🖼️  检测到图片文件，直接使用 OCR: {path.name}")
        try:
            from config.ocr_config import OCR_CONFIG
            
            if OCR_CONFIG.enable_rapidocr:
                ocr_result = extract_with_rapidocr(path, max_pages=1)
                
                if ocr_result.get("error"):
                    result["error"] = ocr_result["error"]
                    result["method"] = "rapidocr_failed"
                    result["confidence"] = 0.0
                else:
                    result["text"] = ocr_result["text"][:limit]
                    result["method"] = "rapidocr"
                    result["confidence"] = ocr_result["confidence"]
            else:
                result["error"] = "RapidOCR 被禁用"
                result["method"] = "disabled"
                result["confidence"] = 0.0
        
        except Exception as e:
            logger.error(f"图片 OCR 失败: {path.name} - {e}")
            result["error"] = str(e)
            result["method"] = "rapidocr_failed"
            result["confidence"] = 0.0
        
        # 🆕 计算质量评分
        quality_score, needs_review = calculate_quality_score(
            result["text"], 
            result["confidence"]
        )
        result["quality_score"] = quality_score
        result["needs_review"] = needs_review
        result["char_count"] = len(result["text"])
        result["processing_time_ms"] = int((time.time() - start_time) * 1000)
        
        # 🆕 保存缓存（如果成功且质量不太差）
        if file_hash and result["text"] and quality_score >= 30:
            try:
                cache.set(
                    file_hash, 
                    result["text"], 
                    result["method"], 
                    result["confidence"], 
                    quality_score
                )
                logger.info(f"💾 OCR 结果已缓存: {path.name}")
            except Exception as e:
                logger.warning(f"保存缓存失败: {e}")
        
        # 记录日志
        logger.info(
            f"图片 OCR 完成: {path.name} | "
            f"方法={result['method']} | "
            f"质量={quality_score} | "
            f"需审查={needs_review} | "
            f"耗时={result['processing_time_ms']}ms"
        )
        
        return result
    
    # PDF 文件 - 使用智能 OCR fallback
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            from config.ocr_config import OCR_CONFIG
            
            # 1. 尝试标准提取
            reader = PdfReader(str(path))
            page_count = len(reader.pages)
            result["page_count"] = page_count
            
            chunks = []
            for page in reader.pages[:10]:  # 前 10 页
                t = page.extract_text() or ""
                if t:
                    chunks.append(t)
                if sum(len(c) for c in chunks) >= limit:
                    break
            
            direct_text = ("\n".join(chunks)).strip()
            
            # 2. 判断是否需要 OCR
            needs_ocr, ocr_reason = should_use_ocr(direct_text, page_count)
            
            if needs_ocr:
                logger.info(f"触发 OCR: {path.name} - {ocr_reason}")
                
                # 3. 判断使用哪种 OCR 方法
                use_vision = (
                    OCR_CONFIG.enable_vision_llm and 
                    page_count <= OCR_CONFIG.important_max_pages and
                    is_important_document(path)
                )
                
                if use_vision:
                    # 使用 Vision LLM（需要调用 core/llm_processor.py）
                    logger.info(f"使用 Vision LLM 处理重要文档: {path.name}")
                    try:
                        from core.llm_processor import analyze_scanned_pdf_with_vision_sync
                        vision_result = analyze_scanned_pdf_with_vision_sync(path)
                        
                        if vision_result.get("error"):
                            # Vision LLM 失败，降级到 RapidOCR
                            logger.warning(f"Vision LLM 失败，降级到 RapidOCR: {vision_result['error']}")
                            ocr_result = extract_with_rapidocr(path, max_pages=OCR_CONFIG.rapidocr_max_pages)
                            result["text"] = ocr_result["text"][:limit]
                            result["method"] = "rapidocr"
                            result["confidence"] = ocr_result["confidence"]
                        else:
                            result["text"] = vision_result["text"][:limit]
                            result["method"] = "vision_llm"
                            result["confidence"] = vision_result["confidence"]
                            
                    except ImportError:
                        logger.warning("Vision LLM 功能未实现，降级到 RapidOCR")
                        ocr_result = extract_with_rapidocr(path, max_pages=OCR_CONFIG.rapidocr_max_pages)
                        result["text"] = ocr_result["text"][:limit]
                        result["method"] = "rapidocr"
                        result["confidence"] = ocr_result["confidence"]
                        
                elif OCR_CONFIG.enable_rapidocr:
                    # 使用 RapidOCR
                    logger.info(f"使用 RapidOCR 处理: {path.name}")
                    ocr_result = extract_with_rapidocr(path, max_pages=OCR_CONFIG.rapidocr_max_pages)
                    
                    if ocr_result.get("error"):
                        # OCR 失败，使用原始文本
                        logger.warning(f"OCR 失败，使用原始提取: {ocr_result['error']}")
                        result["text"] = direct_text[:limit]
                        result["method"] = "direct_fallback"
                        result["confidence"] = 0.3
                        result["error"] = ocr_result["error"]
                    else:
                        result["text"] = ocr_result["text"][:limit]
                        result["method"] = "rapidocr"
                        result["confidence"] = ocr_result["confidence"]
                else:
                    # OCR 被禁用，使用原始文本
                    result["text"] = direct_text[:limit]
                    result["method"] = "direct_fallback"
                    result["confidence"] = 0.3
                    
            else:
                # 不需要 OCR，使用标准提取
                result["text"] = direct_text[:limit]
                result["method"] = "direct"
                result["confidence"] = 0.95
            
            result["char_count"] = len(result["text"])
            
            # 🆕 计算质量评分
            quality_score, needs_review = calculate_quality_score(
                result["text"], 
                result["confidence"]
            )
            result["quality_score"] = quality_score
            result["needs_review"] = needs_review
            
            # 🆕 保存缓存（仅保存 OCR 结果，不缓存 direct 方法）
            if file_hash and result["method"] in ["rapidocr", "vision_llm"] and result["text"]:
                if quality_score >= 30:  # 质量太差不缓存
                    try:
                        cache.set(
                            file_hash, 
                            result["text"], 
                            result["method"], 
                            result["confidence"], 
                            quality_score
                        )
                        logger.info(f"💾 OCR 结果已缓存: {path.name}")
                    except Exception as e:
                        logger.warning(f"保存缓存失败: {e}")
            
        except Exception as e:
            logger.error(f"PDF 处理失败: {path.name} - {e}")
            result["error"] = str(e)
            result["method"] = "direct_fallback"
            result["confidence"] = 0.0
            
            # 🆕 失败情况也计算质量评分
            quality_score, needs_review = calculate_quality_score("", 0.0)
            result["quality_score"] = quality_score
            result["needs_review"] = needs_review
    
    # 其他文件类型保持原有逻辑
    else:
        text = _extract_text_other_formats(path, limit, ext)
        result["text"] = text
        result["method"] = "direct"
        result["confidence"] = 0.9 if text else 0.0
        result["char_count"] = len(text)
        
        # 🆕 计算质量评分
        quality_score, needs_review = calculate_quality_score(
            result["text"], 
            result["confidence"]
        )
        result["quality_score"] = quality_score
        result["needs_review"] = needs_review
    
    # 记录处理时间
    result["processing_time_ms"] = int((time.time() - start_time) * 1000)
    
    # 🆕 记录日志（包含质量评分）
    logger.info(
        f"文本提取完成: {path.name} | "
        f"方法={result['method']} | "
        f"置信度={result['confidence']:.2f} | "
        f"质量={result['quality_score']} | "
        f"需审查={result['needs_review']} | "
        f"字符数={result['char_count']} | "
        f"耗时={result['processing_time_ms']}ms"
    )
    
    return result


def _extract_text_other_formats(path: Path, limit: int, ext: str) -> str:
    """
    提取非 PDF 文件的文本（内部函数）
    """
    # Word 文档
    if ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
            return text.strip()[:limit]
        except Exception:
            return ""

    # PowerPoint 文档
    if ext in [".ppt", ".pptx"]:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
                if sum(len(c) for c in chunks) >= limit:
                    break
            return ("\n".join(chunks)).strip()[:limit]
        except Exception:
            return ""

    # 文本文件
    if ext in [".txt", ".md"]:
        try:
            return path.read_text(encoding="utf-8", errors="ignore").strip()[:limit]
        except Exception:
            return ""

    # 图片 OCR（暂不支持）
    if ext in [".png", ".jpg", ".jpeg", ".webp"]:
        return ""

    return ""


def extract_text_preview(path: Path, limit: int = 1000) -> str:
    """
    提取文件文本预览（兼容旧接口）
    
    Args:
        path: 文件路径
        limit: 最大字符数
    
    Returns:
        文本预览内容
        
    注意：这是兼容旧接口的版本，返回 str。
    新代码建议使用 extract_text_preview_enhanced() 获取更多信息。
    """
    result = extract_text_preview_enhanced(path, limit)
    return result["text"]


# 保持以下代码不变（删除重复的非 PDF 处理代码）
def _old_extract_text_preview_pdf_fallback(path: Path, limit: int) -> str:
    """
    旧版 PDF 提取逻辑（保留作为参考）
    """
    ext = path.suffix.lower()
    
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            chunks = []
            for page in reader.pages[:10]:
                t = page.extract_text() or ""
                if t:
                    chunks.append(t)
                if sum(len(c) for c in chunks) >= limit:
                    break
            return ("\n".join(chunks)).strip()[:limit]
        except Exception:
            return ""


def get_file_size_mb(path: Path) -> float:
    """
    获取文件大小（MB）
    
    Args:
        path: 文件路径
    
    Returns:
        文件大小（MB）
    """
    try:
        return path.stat().st_size / (1024 * 1024)
    except Exception:
        return 0.0


def is_allowed_extension(path: Path, allowed_exts: List[str]) -> bool:
    """
    检查文件扩展名是否在允许列表中
    
    Args:
        path: 文件路径
        allowed_exts: 允许的扩展名列表（如 [".pdf", ".docx"]）
    
    Returns:
        是否允许
    """
    ext = path.suffix.lower()
    return ext in [e.lower() for e in allowed_exts]


# ==================== Step 5: 安全文件移动功能 ====================


def resolve_collision(destination: Path, max_attempts: int = 100) -> Path:
    """
    解决文件名冲突，自动添加数字后缀
    
    当目标文件已存在时，自动在文件名后添加 _1, _2, _3 等后缀，
    直到找到一个不存在的文件名。
    
    Args:
        destination: 目标文件路径
        max_attempts: 最大尝试次数，防止无限循环（默认 100）
    
    Returns:
        不冲突的文件路径
        
    Examples:
        >>> dest = Path("archive/contract.pdf")
        >>> # 如果 contract.pdf 已存在
        >>> result = resolve_collision(dest)
        >>> print(result)  # archive/contract_1.pdf
        
        >>> # 如果 contract.pdf 和 contract_1.pdf 都存在
        >>> result = resolve_collision(dest)
        >>> print(result)  # archive/contract_2.pdf
    
    Raises:
        RuntimeError: 达到最大尝试次数仍未找到可用文件名
    """
    # 如果目标文件不存在，直接返回
    if not destination.exists():
        return destination
    
    # 分离文件名和扩展名
    stem = destination.stem  # 文件名（不含扩展名）
    suffix = destination.suffix  # 扩展名（含 .）
    parent = destination.parent  # 父目录
    
    # 尝试添加后缀 _1, _2, _3...
    for i in range(1, max_attempts + 1):
        new_name = f"{stem}_{i}{suffix}"
        new_path = parent / new_name
        
        if not new_path.exists():
            return new_path
    
    # 达到最大尝试次数仍未找到可用文件名
    raise RuntimeError(
        f"无法解决文件名冲突：已尝试 {max_attempts} 次，"
        f"目标位置 {destination.parent} 可能存在大量同名文件"
    )


def safe_move_file(
    src: Path,
    dst: Path,
    create_dirs: bool = True,
    resolve_conflicts: bool = True
) -> Dict[str, Any]:
    """
    安全地移动文件，包含冲突处理和错误捕获
    
    执行文件移动操作，自动处理目录创建、文件名冲突等问题。
    所有异常都会被捕获并记录在返回字典中，不会向上抛出。
    
    Args:
        src: 源文件路径
        dst: 目标文件路径（可能会因冲突而调整）
        create_dirs: 是否自动创建目标父目录（默认 True）
        resolve_conflicts: 是否自动解决文件名冲突（默认 True）
    
    Returns:
        操作结果字典，包含以下字段：
        - status: str - "success" 或 "failed"
        - src: str - 源文件路径
        - dst: str - 实际移动到的目标路径（可能因冲突而不同于输入的 dst）
        - original_dst: str - 原始请求的目标路径
        - conflict_resolved: bool - 是否解决了文件名冲突
        - error: str - 错误信息（仅在失败时存在）
        
    Examples:
        >>> src = Path("inbox/document.pdf")
        >>> dst = Path("archive/2024/document.pdf")
        >>> result = safe_move_file(src, dst)
        >>> print(result)
        {
            "status": "success",
            "src": "inbox/document.pdf",
            "dst": "archive/2024/document.pdf",
            "original_dst": "archive/2024/document.pdf",
            "conflict_resolved": False
        }
        
        >>> # 如果目标文件已存在
        >>> result = safe_move_file(src, dst)
        >>> print(result)
        {
            "status": "success",
            "src": "inbox/document.pdf",
            "dst": "archive/2024/document_1.pdf",  # 自动添加后缀
            "original_dst": "archive/2024/document.pdf",
            "conflict_resolved": True
        }
    """
    result: Dict[str, Any] = {
        "status": "failed",
        "src": str(src),
        "original_dst": str(dst),
        "dst": str(dst),
        "conflict_resolved": False,
    }
    
    try:
        # 1. 检查源文件是否存在
        if not src.exists():
            result["error"] = f"源文件不存在: {src}"
            return result
        
        if not src.is_file():
            result["error"] = f"源路径不是文件: {src}"
            return result
        
        # 2. 解决文件名冲突
        final_dst = dst
        if resolve_conflicts:
            try:
                final_dst = resolve_collision(dst)
                if final_dst != dst:
                    result["conflict_resolved"] = True
                    result["dst"] = str(final_dst)
            except RuntimeError as e:
                result["error"] = f"无法解决文件名冲突: {e}"
                return result
        
        # 3. 创建目标父目录
        if create_dirs:
            try:
                final_dst.parent.mkdir(parents=True, exist_ok=True)
            except Exception as e:
                result["error"] = f"无法创建目标目录 {final_dst.parent}: {e}"
                return result
        else:
            # 不自动创建目录时，检查目录是否存在
            if not final_dst.parent.exists():
                result["error"] = f"目标目录不存在: {final_dst.parent}"
                return result
        
        # 4. 执行移动操作
        try:
            shutil.move(str(src), str(final_dst))
            result["status"] = "success"
            result["dst"] = str(final_dst)
            # 移除 error 字段（如果之前有的话）
            result.pop("error", None)
            
        except Exception as e:
            result["error"] = f"文件移动失败: {e}"
            return result
        
    except Exception as e:
        # 捕获所有未预期的异常
        result["error"] = f"未知错误: {e}"
    
    return result

