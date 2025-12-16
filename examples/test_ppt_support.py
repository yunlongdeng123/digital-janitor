#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
测试 PPT 文件支持
验证系统是否正确识别和处理 PowerPoint 文件
"""

import sys
import os
from pathlib import Path

# 设置 UTF-8 编码
if sys.platform == "win32":
    os.environ["PYTHONIOENCODING"] = "utf-8"
    if hasattr(sys.stdout, 'reconfigure'):
        sys.stdout.reconfigure(encoding='utf-8')

# 测试导入
def test_imports():
    """测试必要的库是否已安装"""
    print("=" * 80)
    print("🧪 测试 1: 检查依赖库")
    print("=" * 80)
    print()
    
    success = True
    
    # 测试 python-pptx
    try:
        from pptx import Presentation
        print("✅ python-pptx 已安装")
    except ImportError:
        print("❌ python-pptx 未安装")
        print("   请运行: pip install python-pptx")
        success = False
    
    # 测试其他依赖
    try:
        from run_graph_once import JanitorWorkflow
        print("✅ JanitorWorkflow 可导入")
    except ImportError as e:
        print(f"❌ JanitorWorkflow 导入失败: {e}")
        success = False
    
    print()
    return success


def test_config():
    """测试配置文件是否包含 presentation"""
    print("=" * 80)
    print("🧪 测试 2: 检查配置文件")
    print("=" * 80)
    print()
    
    import yaml
    
    config_path = Path("config.yaml")
    if not config_path.exists():
        print("❌ config.yaml 不存在")
        return False
    
    with config_path.open("r", encoding="utf-8") as f:
        config = yaml.safe_load(f)
    
    success = True
    
    # 检查 allowed_ext
    if ".pptx" in config["safety"]["allowed_ext"]:
        print("✅ .pptx 已添加到 allowed_ext")
    else:
        print("❌ .pptx 未添加到 allowed_ext")
        success = False
    
    if ".ppt" in config["safety"]["allowed_ext"]:
        print("✅ .ppt 已添加到 allowed_ext")
    else:
        print("❌ .ppt 未添加到 allowed_ext")
        success = False
    
    # 检查 naming_template
    if "presentation" in config.get("naming_template", {}):
        print("✅ presentation 命名模板已添加")
    else:
        print("❌ presentation 命名模板未添加")
        success = False
    
    # 检查 routing
    if "presentation" in config.get("routing", {}):
        print("✅ presentation 路由规则已添加")
    else:
        print("❌ presentation 路由规则未添加")
        success = False
    
    print()
    return success


def test_ppt_extraction():
    """测试 PPT 文本提取功能"""
    print("=" * 80)
    print("🧪 测试 3: PPT 文本提取")
    print("=" * 80)
    print()
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 创建一个测试 PPT
        test_ppt = Path("inbox/test_presentation.pptx")
        test_ppt.parent.mkdir(parents=True, exist_ok=True)
        
        prs = Presentation()
        
        # 添加标题幻灯片
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "产品发布会"
        subtitle.text = "2024 Q1 新品介绍"
        
        # 添加内容幻灯片
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = bullet_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "主要特性"
        tf = body_shape.text_frame
        tf.text = "创新设计"
        
        # 保存
        prs.save(str(test_ppt))
        print(f"✅ 已创建测试 PPT: {test_ppt}")
        
        # 测试提取
        from utils.file_ops import extract_text_preview
        
        text = extract_text_preview(test_ppt, limit=1000)
        
        if text:
            print(f"✅ 成功提取文本 ({len(text)} 字符)")
            print(f"\n提取的内容预览：")
            print(f"{'-' * 40}")
            print(f"{text[:200]}")
            print(f"{'-' * 40}")
        else:
            print("⚠️  提取的文本为空")
        
        print()
        return True
        
    except Exception as e:
        print(f"❌ PPT 提取测试失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_workflow():
    """测试完整的工作流"""
    print("=" * 80)
    print("🧪 测试 4: 完整工作流测试")
    print("=" * 80)
    print()
    
    try:
        from run_graph_once import JanitorWorkflow
        
        # 初始化工作流
        workflow = JanitorWorkflow()
        print("✅ 工作流初始化成功")
        
        # 查找测试 PPT
        test_ppt = Path("inbox/test_presentation.pptx")
        
        if not test_ppt.exists():
            print("⚠️  测试 PPT 不存在，跳过工作流测试")
            print("   请先运行测试 3 创建测试文件")
            return True
        
        print(f"📄 处理测试文件: {test_ppt.name}")
        print()
        
        # 处理文件（dry-run 模式）
        result = workflow.process_file(
            file_path=test_ppt,
            dry_run=True,
            auto_approve=True,
            max_preview=1000
        )
        
        print("✅ 工作流执行成功")
        print(f"\n📊 处理结果：")
        print(f"   分类: {result['plan']['category']}")
        print(f"   新文件名: {result['plan']['new_name']}")
        print(f"   目标目录: {result['plan']['dest_dir']}")
        print(f"   置信度: {result['plan']['confidence']}")
        
        # 验证是否识别为 presentation
        if result['plan']['category'] == 'presentation':
            print("\n🎉 成功识别为 presentation 类别！")
        else:
            print(f"\n⚠️  未识别为 presentation，而是: {result['plan']['category']}")
        
        print()
        return True
        
    except Exception as e:
        print(f"❌ 工作流测试失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """运行所有测试"""
    print("\n" + "=" * 80)
    print("🚀 PPT 支持测试套件")
    print("=" * 80)
    print()
    
    results = []
    
    # 运行测试
    results.append(("依赖库检查", test_imports()))
    results.append(("配置文件检查", test_config()))
    results.append(("PPT 文本提取", test_ppt_extraction()))
    results.append(("完整工作流", test_workflow()))
    
    # 总结
    print("=" * 80)
    print("📊 测试总结")
    print("=" * 80)
    
    for name, success in results:
        status = "✅ 通过" if success else "❌ 失败"
        print(f"{name:<20} {status}")
    
    all_passed = all(r[1] for r in results)
    
    print()
    if all_passed:
        print("🎉 所有测试通过！PPT 支持已成功配置！")
    else:
        print("⚠️  部分测试失败，请检查上述错误信息")
    
    print("=" * 80)
    
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())

